"""
Pastors' Protocol Central — How It Works
Generates a branded .pptx presentation matching the app's dark-green / gold design.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.dml.color import RGBColor
import copy
from lxml import etree

# ── Palette (matches app CSS) ──────────────────────────────────────────────
BG          = RGBColor(0x04, 0x0B, 0x07)   # near-black green #040b07
BG_PANEL    = RGBColor(0x0D, 0x17, 0x12)   # card background
BG_PANEL2   = RGBColor(0x12, 0x1F, 0x18)   # lighter card
LINE_SOFT   = RGBColor(0x1E, 0x35, 0x28)   # subtle border
GOLD_BRIGHT = RGBColor(0xF1, 0xDD, 0xB0)   # headline gold
GOLD_MAIN   = RGBColor(0xD6, 0xBB, 0x75)   # main gold
GOLD_DIM    = RGBColor(0xA0, 0x8B, 0x52)   # dim gold / eyebrow
TEXT_MAIN   = RGBColor(0xED, 0xF4, 0xEE)   # near-white
TEXT_SOFT   = RGBColor(0xB4, 0xC4, 0xB7)   # muted text
TEXT_MUTED  = RGBColor(0x8E, 0xA1, 0x91)   # faint text
TEXT_FAINT  = RGBColor(0x5A, 0x70, 0x60)   # very faint
GREEN_ACC   = RGBColor(0x22, 0xC5, 0x5E)   # seated / success
AMBER_ACC   = RGBColor(0xF5, 0x9E, 0x0B)   # arrived
RED_ACC     = RGBColor(0xEF, 0x44, 0x44)   # absent
BLUE_SEC    = RGBColor(0x60, 0xA5, 0xFA)   # Middle section
RED_SEC     = RGBColor(0xFC, 0xA5, 0xA5)   # Left section
YELLOW_SEC  = RGBColor(0xFD, 0xE0, 0x47)   # Right section
ORANGE_SEC  = RGBColor(0xFB, 0x92, 0x3C)   # Choir section
PURPLE_SEC  = RGBColor(0xC4, 0xB5, 0xFD)   # VVIP section
SLATE_SEC   = RGBColor(0x94, 0xA3, 0xB8)   # Minister/closed

WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

# ── Fonts ──────────────────────────────────────────────────────────────────
SERIF   = "Garamond"          # closest built-in serif to Libre Bodoni
MONO    = "Courier New"       # Azeret Mono fallback
SANS    = "Calibri"           # Onest fallback

# ── Slide dimensions (widescreen 16:9) ────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)

LOGO_PATH = os.path.join(os.path.dirname(__file__), "..", "public", "logo.png")

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank

# ══════════════════════════════════════════════════════════════════════════
#  Helper utilities
# ══════════════════════════════════════════════════════════════════════════

def add_slide():
    return prs.slides.add_slide(blank_layout)

def rgb(r, g, b):
    return RGBColor(r, g, b)

def fill_slide_bg(slide, color=BG):
    """Fill the slide background with a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    """Add a plain rectangle shape."""
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        l, t, w, h
    )
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def rounded_rect(slide, l, t, w, h, fill_color=BG_PANEL, line_color=LINE_SOFT,
                  line_width=Pt(0.5), adj=0.05):
    """Add a rounded rectangle card."""
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        5,   # MSO_SHAPE_TYPE.ROUNDED_RECTANGLE
        l, t, w, h
    )
    # Set corner radius via adj
    shape.adjustments[0] = adj
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = line_width
    shape.shadow.inherit = False
    return shape

def txbox(slide, text, l, t, w, h,
          font_name=SANS, font_size=Pt(12), bold=False, italic=False,
          color=TEXT_MAIN, align=PP_ALIGN.LEFT, wrap=True,
          line_spacing=None):
    """Add a text box with one paragraph run."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    tf.auto_size = None
    return tb

def multiline_txbox(slide, lines, l, t, w, h,
                     font_name=SANS, font_size=Pt(12), bold=False,
                     color=TEXT_MAIN, align=PP_ALIGN.LEFT,
                     line_spacing=Pt(4)):
    """lines = list of (text, font_name, font_size, bold, color)"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for (txt, fn, fs, b, c) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_after = line_spacing
        run = p.add_run()
        run.text = txt
        run.font.name = fn
        run.font.size = fs
        run.font.bold = b
        run.font.color.rgb = c
    tf.auto_size = None
    return tb

def add_logo(slide, l=Inches(0.35), t=Inches(0.2), w=Inches(1.7), h=Inches(0.55)):
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, l, t, w, h)

def horizontal_rule(slide, l, t, w, color=GOLD_DIM, thickness=Pt(0.5)):
    line = slide.shapes.add_shape(1, l, t, w, Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    line.shadow.inherit = False
    return line

def badge_pill(slide, text, l, t, fill, border, text_color, font_size=Pt(9)):
    w = Inches(1.4)
    h = Inches(0.26)
    shape = slide.shapes.add_shape(5, l, t, w, h)
    shape.adjustments[0] = 0.5
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(0.5)
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = SANS
    run.font.size = font_size
    run.font.bold = True
    run.font.color.rgb = text_color
    return shape

def eyebrow(slide, text, l, t, w=Inches(8)):
    return txbox(slide, text.upper(), l, t, w, Inches(0.28),
                 font_name=MONO, font_size=Pt(8), color=GOLD_DIM,
                 bold=True)

def heading(slide, text, l, t, w=Inches(9), size=Pt(44), color=GOLD_BRIGHT):
    return txbox(slide, text, l, t, w, Inches(1.1),
                 font_name=SERIF, font_size=size, bold=True, color=color)

def body(slide, text, l, t, w=Inches(7), h=Inches(1.4), size=Pt(13), color=TEXT_SOFT):
    return txbox(slide, text, l, t, w, h,
                 font_name=SANS, font_size=size, color=color, wrap=True)

def check_mark(slide, x, y, ok=True):
    sym = "✓" if ok else "—"
    col = GREEN_ACC if ok else TEXT_FAINT
    txbox(slide, sym, x, y, Inches(0.4), Inches(0.3),
          font_size=Pt(14), color=col, bold=True, align=PP_ALIGN.CENTER)

def section_block(slide, label, l, t, w, h, bg, fg, closed=False):
    shape = rounded_rect(slide, l, t, w, h, fill_color=bg,
                         line_color=fg, line_width=Pt(0.8), adj=0.08)
    opacity = 0.5 if closed else 1.0
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label + (" (Closed)" if closed else "")
    run.font.name = SANS
    run.font.size = Pt(8)
    run.font.bold = True
    run.font.color.rgb = fg
    return shape

def stat_card(slide, l, t, value, label, val_color=TEXT_MAIN):
    w, h = Inches(1.6), Inches(0.85)
    rounded_rect(slide, l, t, w, h, fill_color=BG_PANEL, line_color=LINE_SOFT)
    txbox(slide, str(value), l, t + Inches(0.05), w, Inches(0.45),
          font_name=SERIF, font_size=Pt(28), bold=True,
          color=val_color, align=PP_ALIGN.CENTER)
    txbox(slide, label.upper(), l, t + Inches(0.48), w, Inches(0.3),
          font_name=MONO, font_size=Pt(7), color=TEXT_FAINT, align=PP_ALIGN.CENTER)

def seat_dot(slide, l, t, initials, bg, border, fg):
    sz = Inches(0.26)
    shape = rounded_rect(slide, l, t, sz, sz, fill_color=bg,
                          line_color=border, line_width=Pt(0.6), adj=0.15)
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = initials
    run.font.name = SANS
    run.font.size = Pt(6)
    run.font.bold = True
    run.font.color.rgb = fg
    return shape

def empty_dot(slide, l, t):
    sz = Inches(0.26)
    shape = rounded_rect(slide, l, t, sz, sz,
                          fill_color=BG, line_color=LINE_SOFT,
                          line_width=Pt(0.4), adj=0.15)
    return shape

def progress_bar(slide, l, t, w, pct, color, bg=LINE_SOFT):
    h = Inches(0.065)
    # background track
    r1 = slide.shapes.add_shape(1, l, t, w, h)
    r1.fill.solid(); r1.fill.fore_color.rgb = LINE_SOFT
    r1.line.fill.background(); r1.shadow.inherit = False
    # fill
    fill_w = max(Inches(0.05), int(w * pct / 100))
    r2 = slide.shapes.add_shape(1, l, t, fill_w, h)
    r2.fill.solid(); r2.fill.fore_color.rgb = color
    r2.line.fill.background(); r2.shadow.inherit = False

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — Cover
# ══════════════════════════════════════════════════════════════════════════
slide = add_slide()
fill_slide_bg(slide)

# Radial glow accent (simulated with a large soft oval)
glow = slide.shapes.add_shape(9, Inches(-1), Inches(1.5), Inches(5), Inches(4))  # oval
glow.fill.solid(); glow.fill.fore_color.rgb = RGBColor(0x0D, 0x22, 0x14)
glow.line.fill.background(); glow.shadow.inherit = False

# Logo
add_logo(slide, l=Inches(0.55), t=Inches(0.38), w=Inches(2.0), h=Inches(0.65))

# Horizontal rule under logo
horizontal_rule(slide, Inches(0.55), Inches(1.15), Inches(5.8), color=LINE_SOFT)

# Eyebrow
eyebrow(slide, "Pastor's Protocol Central · Presentation", Inches(0.55), Inches(1.28))

# Main title
multiline_txbox(slide,
    [("How It", SERIF, Pt(72), True, GOLD_BRIGHT),
     ("Works",  SERIF, Pt(72), True, GOLD_MAIN)],
    Inches(0.55), Inches(1.65), Inches(6), Inches(2.6),
    align=PP_ALIGN.LEFT, line_spacing=Pt(2))

# Subtitle
body(slide,
     "A guided tour of the dignitary seating management system\n"
     "— from signing in to the live control centre.",
     Inches(0.55), Inches(4.05), Inches(5.6), Inches(0.9),
     size=Pt(13), color=TEXT_SOFT)

# Badges row
bx = Inches(0.55)
badge_pill(slide, "Conferences",      bx,              Inches(5.12), BG_PANEL, GOLD_DIM,  GOLD_MAIN)
badge_pill(slide, "Live Attendance",  bx+Inches(1.52), Inches(5.12), BG_PANEL, LINE_SOFT, GREEN_ACC)
badge_pill(slide, "Seating Maps",     bx+Inches(3.04), Inches(5.12), BG_PANEL, LINE_SOFT, TEXT_SOFT)
badge_pill(slide, "Role-based Access",bx+Inches(4.56), Inches(5.12), BG_PANEL, LINE_SOFT, TEXT_SOFT)

# Right-side stat column
sx = Inches(8.8)
horizontal_rule(slide, sx, Inches(1.2), Inches(0.01), color=GOLD_DIM, thickness=Pt(0.5))
# vertical rule
vr = slide.shapes.add_shape(1, sx + Inches(0.3), Inches(1.2), Pt(1), Inches(5.2))
vr.fill.solid(); vr.fill.fore_color.rgb = RGBColor(0x1A, 0x35, 0x24)
vr.line.fill.background(); vr.shadow.inherit = False

for i, (num, lbl) in enumerate([("12", "App Screens"), ("3", "Access Roles"), ("∞", "Conferences")]):
    sy = Inches(1.5) + i * Inches(1.55)
    txbox(slide, num, sx+Inches(0.55), sy, Inches(2.2), Inches(0.85),
          font_name=SERIF, font_size=Pt(52), bold=True,
          color=GOLD_BRIGHT, align=PP_ALIGN.CENTER)
    txbox(slide, lbl.upper(), sx+Inches(0.55), sy+Inches(0.78), Inches(2.2), Inches(0.3),
          font_name=MONO, font_size=Pt(8), color=TEXT_FAINT, align=PP_ALIGN.CENTER)
    if i < 2:
        hr = slide.shapes.add_shape(1, sx+Inches(0.7), sy+Inches(1.18), Inches(1.8), Pt(1))
        hr.fill.solid(); hr.fill.fore_color.rgb = LINE_SOFT
        hr.line.fill.background(); hr.shadow.inherit = False

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — The Challenge
# ══════════════════════════════════════════════════════════════════════════
slide = add_slide()
fill_slide_bg(slide)
add_logo(slide)

eyebrow(slide, "Why We Built It", Inches(0.55), Inches(0.82))
heading(slide, "The Challenge of Large Conferences",
        Inches(0.55), Inches(1.08), size=Pt(38))

body(slide,
     "Managing where hundreds of pastors and dignitaries sit across multiple sessions is complex, "
     "error-prone, and stressful without the right tools.",
     Inches(0.55), Inches(2.02), Inches(8.5), Inches(0.7), size=Pt(12))

# 4 problem cards in a 2×2 grid
problems = [
    ("📋", "Spreadsheet Chaos",
     "Paper lists and spreadsheets get lost, go out of date, and can't show real-time arrival status."),
    ("🪑", "No Visual Overview",
     "Without a floor plan view, protocol officers can't quickly see which seats are taken or empty."),
    ("👥", "Team Coordination",
     "Multiple officers working from different lists leads to double-bookings and missed guests."),
    ("🔄", "Repeated Setup",
     "Each conference session means starting over — re-entering the same pastors, re-assigning seats."),
]
cols = [Inches(0.55), Inches(6.85)]
rows = [Inches(2.88), Inches(5.02)]
cw, ch = Inches(5.9), Inches(1.85)

for idx, (icon, title, desc) in enumerate(problems):
    cx = cols[idx % 2]
    cy = rows[idx // 2]
    rounded_rect(slide, cx, cy, cw, ch, fill_color=BG_PANEL, line_color=LINE_SOFT)
    txbox(slide, icon, cx+Inches(0.18), cy+Inches(0.14), Inches(0.5), Inches(0.45), font_size=Pt(20))
    txbox(slide, title, cx+Inches(0.7), cy+Inches(0.14), cw-Inches(0.9), Inches(0.4),
          font_name=SERIF, font_size=Pt(16), bold=True, color=TEXT_MAIN)
    txbox(slide, desc, cx+Inches(0.18), cy+Inches(0.62), cw-Inches(0.36), Inches(1.1),
          font_name=SANS, font_size=Pt(11), color=TEXT_MUTED, wrap=True)

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — The Solution
# ══════════════════════════════════════════════════════════════════════════
slide = add_slide()
fill_slide_bg(slide)
add_logo(slide)

eyebrow(slide, "The Solution", Inches(0.55), Inches(0.82))
heading(slide, "One System, Every Event",
        Inches(0.55), Inches(1.08), size=Pt(42))

body(slide,
     "Pastors' Protocol Central brings every part of dignitary management into a single, beautiful "
     "web app your whole protocol team can access together.",
     Inches(0.55), Inches(2.0), Inches(9), Inches(0.65), size=Pt(12))

# 4-step flow cards
steps = [
    ("🏟️", "Set Up", "Create the event, pick the auditorium, add sessions."),
    ("👤", "Build Roster", "Add dignitaries to a shared directory, assign to sessions."),
    ("🗺️", "Assign Seats", "Click the floor plan to place each pastor's seat."),
    ("📡", "Track Live", "Mark arrivals and seating in real time from control centre."),
]

sw = Inches(2.75); sh = Inches(3.1)
gap = Inches(0.25)
start_x = Inches(0.55)
sy_start = Inches(2.88)

for i, (icon, title, desc) in enumerate(steps):
    sx = start_x + i * (sw + gap)
    rounded_rect(slide, sx, sy_start, sw, sh, fill_color=BG_PANEL2, line_color=LINE_SOFT)
    # Number
    txbox(slide, str(i+1), sx+Inches(0.18), sy_start+Inches(0.15), Inches(0.5), Inches(0.55),
          font_name=SERIF, font_size=Pt(32), bold=True, color=GOLD_DIM)
    # Icon
    txbox(slide, icon, sx+sw-Inches(0.6), sy_start+Inches(0.18), Inches(0.5), Inches(0.5),
          font_size=Pt(22))
    # Title
    txbox(slide, title, sx+Inches(0.18), sy_start+Inches(0.82), sw-Inches(0.36), Inches(0.5),
          font_name=SERIF, font_size=Pt(18), bold=True, color=TEXT_MAIN)
    # Desc
    txbox(slide, desc, sx+Inches(0.18), sy_start+Inches(1.4), sw-Inches(0.36), Inches(1.5),
          font_name=SANS, font_size=Pt(11), color=TEXT_MUTED, wrap=True)
    # Arrow between cards
    if i < 3:
        ax = sx + sw + Inches(0.04)
        ay = sy_start + sh/2 - Inches(0.15)
        txbox(slide, "›", ax, ay, gap+Inches(0.05), Inches(0.35),
              font_size=Pt(22), color=TEXT_FAINT, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — Sign In & Roles
# ══════════════════════════════════════════════════════════════════════════
slide = add_slide()
fill_slide_bg(slide)
add_logo(slide)

eyebrow(slide, "Step 1", Inches(0.55), Inches(0.82))
heading(slide, "Sign In", Inches(0.55), Inches(1.08), size=Pt(42))

body(slide,
     "Everyone gets their own secure login. Access is tailored to what your role needs.",
     Inches(0.55), Inches(2.0), Inches(4.2), Inches(0.65), size=Pt(12))

# Auth card mockup
ax, ay, aw, ah = Inches(0.55), Inches(2.88), Inches(3.4), Inches(4.1)
rounded_rect(slide, ax, ay, aw, ah, fill_color=BG_PANEL, line_color=LINE_SOFT)

txbox(slide, "Pastors' Protocol Central",
      ax+Inches(0.1), ay+Inches(0.18), aw-Inches(0.2), Inches(0.4),
      font_name=SERIF, font_size=Pt(13), bold=True, color=GOLD_MAIN, align=PP_ALIGN.CENTER)
txbox(slide, "DIGNITARY MANAGEMENT",
      ax+Inches(0.1), ay+Inches(0.56), aw-Inches(0.2), Inches(0.26),
      font_name=MONO, font_size=Pt(7), color=TEXT_FAINT, align=PP_ALIGN.CENTER)

# Tabs
tab_y = ay + Inches(0.9)
rounded_rect(slide, ax+Inches(0.14), tab_y, (aw-Inches(0.28))/2, Inches(0.35),
             fill_color=BG_PANEL2, line_color=GOLD_DIM)
txbox(slide, "Sign In", ax+Inches(0.14), tab_y, (aw-Inches(0.28))/2, Inches(0.35),
      font_name=SANS, font_size=Pt(10), bold=True, color=GOLD_MAIN, align=PP_ALIGN.CENTER)
rounded_rect(slide, ax+Inches(0.14)+(aw-Inches(0.28))/2, tab_y,
             (aw-Inches(0.28))/2, Inches(0.35), fill_color=BG, line_color=LINE_SOFT)
txbox(slide, "Register", ax+Inches(0.14)+(aw-Inches(0.28))/2, tab_y,
      (aw-Inches(0.28))/2, Inches(0.35),
      font_name=SANS, font_size=Pt(10), color=TEXT_FAINT, align=PP_ALIGN.CENTER)

# Fields
for i, lbl in enumerate(["📧  email@church.org", "🔒  ••••••••••"]):
    fy = ay + Inches(1.42) + i*Inches(0.6)
    rounded_rect(slide, ax+Inches(0.14), fy, aw-Inches(0.28), Inches(0.46),
                 fill_color=BG, line_color=LINE_SOFT)
    txbox(slide, lbl, ax+Inches(0.28), fy+Inches(0.08), aw-Inches(0.56), Inches(0.32),
          font_name=SANS, font_size=Pt(10), color=TEXT_MUTED)

# Sign In button
btn_y = ay + Inches(2.7)
shape = slide.shapes.add_shape(5, ax+Inches(0.14), btn_y, aw-Inches(0.28), Inches(0.46))
shape.adjustments[0] = 0.5
shape.fill.solid(); shape.fill.fore_color.rgb = GOLD_MAIN
shape.line.fill.background(); shape.shadow.inherit = False
tf = shape.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run = p.add_run(); run.text = "SIGN IN"
run.font.name = SANS; run.font.size = Pt(10); run.font.bold = True
run.font.color.rgb = RGBColor(0x09, 0x11, 0x0C)

txbox(slide, "New users can register and will receive Protocol Member access.",
      ax+Inches(0.1), ay+Inches(3.3), aw-Inches(0.2), Inches(0.55),
      font_name=SANS, font_size=Pt(9), color=TEXT_FAINT, align=PP_ALIGN.CENTER, wrap=True)

# Roles table
tx, ty = Inches(4.2), Inches(2.78)
tw, th = Inches(8.65), Inches(4.2)
rounded_rect(slide, tx, ty, tw, th, fill_color=BG_PANEL, line_color=LINE_SOFT)

# Table header
txbox(slide, "WHO CAN DO WHAT", tx+Inches(0.2), ty+Inches(0.15), tw-Inches(0.4), Inches(0.3),
      font_name=MONO, font_size=Pt(8), color=GOLD_DIM, bold=True)

cols_labels = ["Role", "Conferences", "Add Pastors", "Edit Seating", "View"]
col_x = [tx+Inches(0.18), tx+Inches(2.4), tx+Inches(4.0), tx+Inches(5.6), tx+Inches(7.5)]
header_y = ty + Inches(0.55)
for j, lbl in enumerate(cols_labels):
    txbox(slide, lbl.upper(), col_x[j], header_y, Inches(1.8), Inches(0.28),
          font_name=MONO, font_size=Pt(7), color=TEXT_FAINT, bold=True)

horizontal_rule(slide, tx+Inches(0.18), ty+Inches(0.88), tw-Inches(0.36), color=LINE_SOFT)

roles = [
    ("Administrator", "Full system control", True,  True,  True,  True),
    ("Editor",        "Manages content",     True,  True,  True,  True),
    ("Protocol Member","View & update status",False, False, False, True),
]

for r, (role, sub, c1, c2, c3, c4) in enumerate(roles):
    ry = ty + Inches(1.0) + r * Inches(0.96)
    if r % 2 == 0:
        rounded_rect(slide, tx+Inches(0.08), ry-Inches(0.06), tw-Inches(0.16), Inches(0.9),
                     fill_color=RGBColor(0x10, 0x22, 0x18), line_color=RGBColor(0x10,0x22,0x18))
    txbox(slide, role, col_x[0], ry, Inches(2.0), Inches(0.36),
          font_name=SANS, font_size=Pt(12), bold=True, color=TEXT_MAIN)
    txbox(slide, sub, col_x[0], ry+Inches(0.34), Inches(2.0), Inches(0.3),
          font_name=SANS, font_size=Pt(9), color=TEXT_FAINT)
    for ci, ok in enumerate([c1, c2, c3, c4]):
        cx2 = col_x[ci+1]
        sym = "✓" if ok else "—"
        col2 = GREEN_ACC if ok else TEXT_FAINT
        txbox(slide, sym, cx2, ry+Inches(0.12), Inches(1.6), Inches(0.36),
              font_name=SANS, font_size=Pt(14), bold=True,
              color=col2, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — Conferences & Sessions
# ══════════════════════════════════════════════════════════════════════════
slide = add_slide()
fill_slide_bg(slide)
add_logo(slide)

eyebrow(slide, "Step 2", Inches(0.55), Inches(0.82))
heading(slide, "Conferences & Sessions", Inches(0.55), Inches(1.08), size=Pt(40))

body(slide,
     "Each church event is a Conference. Within it you create Sessions — one per service or day — "
     "each with its own independent seating plan.",
     Inches(0.55), Inches(2.0), Inches(9), Inches(0.7), size=Pt(12))

# Conference cards
for i, (title, meta, tag) in enumerate([
    ("General Council 2025", "12 Jun 2025 · 09:00 · National Auditorium, Accra", "GLT IFE Auditorium"),
    ("Camp Meeting 2025",    "20 Aug 2025 · 18:00 · GLT Camp Ground",             "GLT Camp Auditorium"),
]):
    cx = Inches(0.55) + i * Inches(4.55)
    cy = Inches(2.9)
    cw2, ch2 = Inches(4.2), Inches(1.85)
    rounded_rect(slide, cx, cy, cw2, ch2, fill_color=BG_PANEL, line_color=LINE_SOFT)
    txbox(slide, "CONFERENCE", cx+Inches(0.18), cy+Inches(0.14), cw2-Inches(0.36), Inches(0.25),
          font_name=MONO, font_size=Pt(7), color=TEXT_FAINT)
    txbox(slide, title, cx+Inches(0.18), cy+Inches(0.36), cw2-Inches(0.36), Inches(0.46),
          font_name=SERIF, font_size=Pt(18), bold=True, color=TEXT_MAIN)
    txbox(slide, meta, cx+Inches(0.18), cy+Inches(0.82), cw2-Inches(0.36), Inches(0.36),
          font_name=SANS, font_size=Pt(9), color=TEXT_FAINT)
    badge_pill(slide, tag, cx+Inches(0.18), cy+Inches(1.26),
               BG_PANEL2, LINE_SOFT, TEXT_SOFT, font_size=Pt(8))

# Sessions list card
sx2, sy2, sw2, sh2 = Inches(9.35), Inches(2.9), Inches(3.8), Inches(3.85)
rounded_rect(slide, sx2, sy2, sw2, sh2, fill_color=BG_PANEL, line_color=LINE_SOFT)

txbox(slide, "SESSIONS — GENERAL COUNCIL 2025",
      sx2+Inches(0.18), sy2+Inches(0.18), sw2-Inches(0.36), Inches(0.28),
      font_name=MONO, font_size=Pt(7), color=TEXT_FAINT)

sessions_data = [
    ("Opening Service",       "12 Jun 2025 · 09:00"),
    ("Afternoon Session",     "12 Jun 2025 · 14:00"),
    ("Evening Gala Service",  "12 Jun 2025 · 18:30"),
]
for si, (sname, smeta) in enumerate(sessions_data):
    srow_y = sy2 + Inches(0.6) + si * Inches(1.02)
    if si % 2 == 0:
        rounded_rect(slide, sx2+Inches(0.1), srow_y, sw2-Inches(0.2), Inches(0.96),
                     fill_color=RGBColor(0x10,0x22,0x18), line_color=RGBColor(0x10,0x22,0x18))
    txbox(slide, sname, sx2+Inches(0.22), srow_y+Inches(0.12), sw2-Inches(0.44), Inches(0.42),
          font_name=SERIF, font_size=Pt(15), bold=True, color=TEXT_MAIN)
    txbox(slide, smeta, sx2+Inches(0.22), srow_y+Inches(0.5), sw2-Inches(0.44), Inches(0.3),
          font_name=SANS, font_size=Pt(10), color=TEXT_FAINT)

txbox(slide, "💡 Each session keeps its own seating plan — great for conferences where pastors move between services.",
      Inches(0.55), Inches(6.72), Inches(8.5), Inches(0.5),
      font_name=SANS, font_size=Pt(10), color=TEXT_FAINT, wrap=True)

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — Dignitary Directory
# ══════════════════════════════════════════════════════════════════════════
slide = add_slide()
fill_slide_bg(slide)
add_logo(slide)

eyebrow(slide, "Step 3", Inches(0.55), Inches(0.82))
heading(slide, "Dignitary Directory", Inches(0.55), Inches(1.08), size=Pt(42))

body(slide,
     "Build a central directory of all pastors and dignitaries. Each record holds their photo, title, "
     "church, branch, and protocol notes. Add them once — reuse them across sessions forever.",
     Inches(0.55), Inches(2.0), Inches(9), Inches(0.7), size=Pt(12))

# Directory cards in a 3×2 grid
dir_entries = [
    ("AM", "Archbishop Mensah",    "General Overseer",   "GLT International · HQ",      RGBColor(0xC9,0x70,0x3C), "Photo ✓",    GOLD_MAIN),
    ("PO", "Pastor Oti-Boateng",   "Senior Pastor",      "GLT IFE · North Campus",       RGBColor(0x42,0x78,0xB4), "Notes added",TEXT_SOFT),
    ("BK", "Bishop Kwarteng",      "District Bishop",    "GLT International · East",     RGBColor(0x42,0xA3,0x67), "Photo ✓",    GOLD_MAIN),
    ("PA", "Pastor Asamoah",       "Associate Pastor",   "GLT International · West",     RGBColor(0xA3,0x42,0x42), "",           None),
    ("HE", "H.E. Rev. Egyir",      "Presiding Bishop",   "GLT Camp · Director",          RGBColor(0x78,0x42,0xA3), "Notes added",TEXT_SOFT),
    ("+",  "Add New Dignitary",    "",                   "",                             LINE_SOFT, "",            None),
]
dw, dh = Inches(3.85), Inches(1.82)
gap_d = Inches(0.25)
start_dx = Inches(0.55)
start_dy = Inches(2.9)

for idx, (init, name, title2, church, av_col, badge_txt, badge_col) in enumerate(dir_entries):
    col = idx % 3; row = idx // 3
    dx = start_dx + col * (dw + gap_d)
    dy = start_dy + row * (dh + Inches(0.2))
    if init == "+":
        rounded_rect(slide, dx, dy, dw, dh, fill_color=BG, line_color=LINE_SOFT,
                     line_width=Pt(0.5), adj=0.06)
        txbox(slide, "+ Add New Dignitary", dx, dy+Inches(0.7), dw, Inches(0.45),
              font_name=SANS, font_size=Pt(12), color=TEXT_FAINT, align=PP_ALIGN.CENTER)
        continue
    rounded_rect(slide, dx, dy, dw, dh, fill_color=BG_PANEL, line_color=LINE_SOFT)
    # Avatar circle
    ov = slide.shapes.add_shape(9, dx+Inches(0.18), dy+Inches(0.22), Inches(0.6), Inches(0.6))
    ov.fill.solid()
    r2, g2, b2 = av_col[0], av_col[1], av_col[2]
    ov.fill.fore_color.rgb = RGBColor(r2//4, g2//4, b2//4)
    ov.line.color.rgb = av_col; ov.line.width = Pt(1); ov.shadow.inherit = False
    txbox(slide, init, dx+Inches(0.18), dy+Inches(0.24), Inches(0.6), Inches(0.45),
          font_name=SERIF, font_size=Pt(13), bold=True, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
    txbox(slide, name, dx+Inches(0.92), dy+Inches(0.18), dw-Inches(1.1), Inches(0.36),
          font_name=SANS, font_size=Pt(12), bold=True, color=TEXT_MAIN)
    txbox(slide, title2, dx+Inches(0.92), dy+Inches(0.52), dw-Inches(1.1), Inches(0.3),
          font_name=SANS, font_size=Pt(10), color=TEXT_MUTED)
    txbox(slide, church, dx+Inches(0.18), dy+Inches(0.96), dw-Inches(0.36), Inches(0.28),
          font_name=SANS, font_size=Pt(9), color=TEXT_FAINT)
    if badge_txt and badge_col:
        badge_pill(slide, badge_txt, dx+Inches(0.18), dy+Inches(1.36),
                   BG_PANEL2, LINE_SOFT, badge_col, font_size=Pt(8))

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — Venue Map & Seat Grid
# ══════════════════════════════════════════════════════════════════════════
slide = add_slide()
fill_slide_bg(slide)
add_logo(slide)

eyebrow(slide, "Step 4", Inches(0.55), Inches(0.82))
heading(slide, "Interactive Seating Map", Inches(0.55), Inches(1.08), size=Pt(40))

body(slide,
     "Click any section on the visual floor plan to open its seat grid. Click an empty seat to assign "
     "a dignitary — or click an occupied seat to view their profile.",
     Inches(0.55), Inches(2.0), Inches(9), Inches(0.65), size=Pt(12))

# ── Floor plan (left card) ──
fp_x, fp_y = Inches(0.55), Inches(2.85)
fp_w, fp_h = Inches(5.5), Inches(4.2)
rounded_rect(slide, fp_x, fp_y, fp_w, fp_h, fill_color=BG_PANEL, line_color=LINE_SOFT)
txbox(slide, "AUDITORIUM FLOOR PLAN",
      fp_x+Inches(0.18), fp_y+Inches(0.18), fp_w-Inches(0.36), Inches(0.26),
      font_name=MONO, font_size=Pt(7), color=TEXT_FAINT)

# Altar
section_block(slide, "ALTAR",
              fp_x+Inches(1.5), fp_y+Inches(0.56),
              Inches(2.3), Inches(0.62),
              RGBColor(0x14,0x52,0x22), GREEN_ACC)

# VVIP
section_block(slide, "VVIP",
              fp_x+Inches(2.05), fp_y+Inches(1.3),
              Inches(1.2), Inches(0.5),
              RGBColor(0x2A,0x18,0x45), PURPLE_SEC)

# Row of sections
row_y = fp_y + Inches(2.0)
row_defs = [
    ("CHOIR", RGBColor(0x40,0x20,0x10), ORANGE_SEC, False),
    ("LEFT",  RGBColor(0x40,0x0A,0x0A), RED_SEC,    False),
    ("MIDDLE ◀", RGBColor(0x0A,0x18,0x3C), BLUE_SEC, False),
    ("RIGHT", RGBColor(0x38,0x2A,0x06), YELLOW_SEC, False),
    ("MINISTER", RGBColor(0x1A,0x22,0x2A), SLATE_SEC, True),
]
col_w2 = (fp_w - Inches(0.36)) / len(row_defs) - Inches(0.06)
for ri, (lbl, bg2, fg2, closed) in enumerate(row_defs):
    rx2 = fp_x + Inches(0.18) + ri * (col_w2 + Inches(0.06))
    shape = section_block(slide, lbl, rx2, row_y, col_w2, Inches(1.1), bg2, fg2, closed=closed)

# Entrances
txbox(slide, "ENTRANCE", fp_x+Inches(0.18), fp_y+fp_h-Inches(0.36), Inches(1.2), Inches(0.26),
      font_name=MONO, font_size=Pt(7), color=TEXT_FAINT)
txbox(slide, "ENTRANCE", fp_x+fp_w-Inches(1.5), fp_y+fp_h-Inches(0.36), Inches(1.2), Inches(0.26),
      font_name=MONO, font_size=Pt(7), color=TEXT_FAINT)

# ── Seat Grid (right card) ──
sg_x, sg_y = Inches(6.35), Inches(2.85)
sg_w, sg_h = Inches(6.7), Inches(4.2)
rounded_rect(slide, sg_x, sg_y, sg_w, sg_h, fill_color=BG_PANEL, line_color=LINE_SOFT)

txbox(slide, "Middle Section   (4 rows × 6 seats)",
      sg_x+Inches(0.18), sg_y+Inches(0.18), sg_w-Inches(0.36), Inches(0.36),
      font_name=SERIF, font_size=Pt(14), bold=True, color=TEXT_MAIN)

# Legend
leg_data = [("Pending","#94A3B8"), ("Arrived","#F59E0B"), ("Seated","#22C55E"), ("Absent","#EF4444")]
for li, (lbl, col_hex) in enumerate(leg_data):
    lx = sg_x + Inches(0.18) + li * Inches(1.55)
    dot = slide.shapes.add_shape(9, lx, sg_y+Inches(0.66), Inches(0.13), Inches(0.13))
    c = RGBColor(int(col_hex[1:3],16), int(col_hex[3:5],16), int(col_hex[5:7],16))
    dot.fill.solid(); dot.fill.fore_color.rgb = c
    dot.line.fill.background(); dot.shadow.inherit = False
    txbox(slide, lbl, lx+Inches(0.18), sg_y+Inches(0.6), Inches(1.3), Inches(0.26),
          font_name=SANS, font_size=Pt(9), color=TEXT_MUTED)

# Grid data  (initials, type: "seated","arrived","pending","empty")
grid = [
    [("AM","seated"),("PO","arrived"),("BK","pending"),("","empty"),("","empty"),("PA","seated")],
    [("HE","arrived"),("","empty"),("","empty"),("RA","pending"),("","empty"),("","empty")],
    [("","empty"),("","empty"),("YA","seated"),("","empty"),("NK","arrived"),("","empty")],
    [("","empty"),("","empty"),("","empty"),("","empty"),("","empty"),("","empty")],
]
type_colors = {
    "seated":  (RGBColor(0x05,0x29,0x16), GREEN_ACC,  RGBColor(0x22,0xC5,0x5E)),
    "arrived": (RGBColor(0x35,0x1A,0x03), AMBER_ACC,  RGBColor(0xF5,0x9E,0x0B)),
    "pending": (RGBColor(0x1E,0x29,0x3B), LINE_SOFT,  RGBColor(0x94,0xA3,0xB8)),
    "empty":   (BG, LINE_SOFT, BG),
}
dot_sz = Inches(0.3)
dot_gap = Inches(0.06)
grid_start_x = sg_x + Inches(0.22)
grid_start_y = sg_y + Inches(0.98)

for ri2, row2 in enumerate(grid):
    for ci2, (init2, typ) in enumerate(row2):
        dx2 = grid_start_x + ci2 * (dot_sz + dot_gap)
        dy2 = grid_start_y + ri2 * (dot_sz + dot_gap)
        bg3, border3, fg3 = type_colors[typ]
        if typ == "empty":
            s2 = rounded_rect(slide, dx2, dy2, dot_sz, dot_sz,
                              fill_color=BG, line_color=LINE_SOFT,
                              line_width=Pt(0.4), adj=0.15)
        else:
            s2 = rounded_rect(slide, dx2, dy2, dot_sz, dot_sz,
                              fill_color=bg3, line_color=border3,
                              line_width=Pt(0.6), adj=0.15)
            tf3 = s2.text_frame
            p3 = tf3.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
            r3 = p3.add_run(); r3.text = init2
            r3.font.name = SANS; r3.font.size = Pt(6)
            r3.font.bold = True; r3.font.color.rgb = fg3

txbox(slide, "Click an empty seat to assign · Click a filled seat to view profile",
      sg_x+Inches(0.18), sg_y+sg_h-Inches(0.4), sg_w-Inches(0.36), Inches(0.3),
      font_name=SANS, font_size=Pt(9), color=TEXT_FAINT)

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — Live Attendance Tracking
# ══════════════════════════════════════════════════════════════════════════
slide = add_slide()
fill_slide_bg(slide)
add_logo(slide)

eyebrow(slide, "Step 5", Inches(0.55), Inches(0.82))
heading(slide, "Live Attendance Tracking", Inches(0.55), Inches(1.08), size=Pt(40))

body(slide,
     "As the service begins, protocol officers update each dignitary's status with one tap. "
     "Every change is instantly visible to the whole team.",
     Inches(0.55), Inches(2.0), Inches(9), Inches(0.65), size=Pt(12))

# Stats bar
stat_data = [("18",TEXT_MAIN,"In Session"),("4",SLATE_SEC,"Pending"),
             ("7",AMBER_ACC,"Arrived"),("6",GREEN_ACC,"Seated"),("1",RED_ACC,"Absent")]
for si2, (val, col3, lbl2) in enumerate(stat_data):
    sx3 = Inches(0.55) + si2 * Inches(2.5)
    stat_card(slide, sx3, Inches(2.84), val, lbl2, col3)

# Attendee list
att_data = [
    ("AM","Archbishop Mensah","General Overseer · GLT International",
     "Seated", GREEN_ACC, "Middle", BLUE_SEC, "R1 / S1"),
    ("PO","Pastor Oti-Boateng","Senior Pastor · GLT IFE North",
     "Arrived", AMBER_ACC, "Middle", BLUE_SEC, "R1 / S2"),
    ("HE","H.E. Rev. Egyir","Presiding Bishop · GLT Camp",
     "Absent", RED_ACC, "Left", RED_SEC, "R2 / S1"),
]
list_y = Inches(3.88)
list_h_per = Inches(0.98)
for ai, (init3, name3, sub3, status3, s_col, sec3, sec_col, seat3) in enumerate(att_data):
    ay2 = list_y + ai * list_h_per
    aw2 = Inches(12.25)
    bg4 = RGBColor(0x10,0x22,0x18) if ai % 2 == 0 else BG_PANEL
    rounded_rect(slide, Inches(0.55), ay2, aw2, list_h_per - Inches(0.06),
                 fill_color=bg4, line_color=LINE_SOFT)

    # Avatar
    ov2 = slide.shapes.add_shape(9, Inches(0.75), ay2+Inches(0.2), Inches(0.6), Inches(0.6))
    ov2.fill.solid(); ov2.fill.fore_color.rgb = RGBColor(0x14,0x22,0x18)
    ov2.line.color.rgb = s_col; ov2.line.width = Pt(1); ov2.shadow.inherit = False
    txbox(slide, init3, Inches(0.75), ay2+Inches(0.2), Inches(0.6), Inches(0.48),
          font_name=SERIF, font_size=Pt(12), bold=True, color=TEXT_MAIN, align=PP_ALIGN.CENTER)

    txbox(slide, name3, Inches(1.5), ay2+Inches(0.1), Inches(3.8), Inches(0.4),
          font_name=SANS, font_size=Pt(12), bold=True, color=TEXT_MAIN)
    txbox(slide, sub3, Inches(1.5), ay2+Inches(0.48), Inches(3.8), Inches(0.3),
          font_name=SANS, font_size=Pt(9.5), color=TEXT_MUTED)

    # Status badge
    badge_pill(slide, f"● {status3}", Inches(5.6), ay2+Inches(0.32),
               RGBColor(0x0A,0x16,0x0E), s_col, s_col, font_size=Pt(9))

    # Section badge
    badge_pill(slide, sec3, Inches(7.15), ay2+Inches(0.32),
               RGBColor(0x06,0x10,0x22), sec_col, sec_col, font_size=Pt(9))

    # Seat ref
    txbox(slide, seat3, Inches(8.75), ay2+Inches(0.32), Inches(1.5), Inches(0.3),
          font_name=MONO, font_size=Pt(9), color=TEXT_FAINT)

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — Control Centre
# ══════════════════════════════════════════════════════════════════════════
slide = add_slide()
fill_slide_bg(slide)
add_logo(slide)

eyebrow(slide, "Step 6 · Live Conference Command", Inches(0.55), Inches(0.82))
heading(slide, "Control Centre", Inches(0.55), Inches(1.08), size=Pt(42))

body(slide,
     "The operations hub — attendance metrics, live alerts, activity feed, and seating coverage for "
     "the entire conference on one screen.",
     Inches(0.55), Inches(2.0), Inches(9), Inches(0.6), size=Pt(12))

# Metric cards row
cc_metrics = [
    ("👥", "48", "Expected"),
    ("✅", "31", "Arrived"),
    ("🪑", "26", "Seated"),
    ("🛡️", "88%", "Protocol Coverage"),
    ("⚠️", "2",  "Needs Attention"),
]
cc_colors = [TEXT_MAIN, AMBER_ACC, GREEN_ACC, GOLD_MAIN, RED_ACC]
cc_mw = Inches(2.35)
for mi, ((icon4, val4, lbl4), mc) in enumerate(zip(cc_metrics, cc_colors)):
    mx = Inches(0.55) + mi * (cc_mw + Inches(0.14))
    my = Inches(2.82)
    rounded_rect(slide, mx, my, cc_mw, Inches(1.1), fill_color=BG_PANEL, line_color=LINE_SOFT)
    txbox(slide, icon4, mx+Inches(0.14), my+Inches(0.14), Inches(0.5), Inches(0.42), font_size=Pt(18))
    txbox(slide, val4, mx+Inches(0.64), my+Inches(0.12), cc_mw-Inches(0.78), Inches(0.48),
          font_name=SERIF, font_size=Pt(26), bold=True, color=mc)
    txbox(slide, lbl4.upper(), mx+Inches(0.14), my+Inches(0.72), cc_mw-Inches(0.28), Inches(0.28),
          font_name=MONO, font_size=Pt(7), color=TEXT_FAINT)

# Three panels below
panels = [
    ("Attention Needed", "Conflicts, gaps, and coverage risks"),
    ("Live Activity", "Latest movements across the event"),
    ("Session Seating", "Section occupancy at a glance"),
]
pw = Inches(4.1); ph = Inches(2.7)
pgap = Inches(0.24)
for pi, (ptitle, psub) in enumerate(panels):
    px = Inches(0.55) + pi * (pw + pgap)
    py = Inches(4.08)
    rounded_rect(slide, px, py, pw, ph, fill_color=BG_PANEL, line_color=LINE_SOFT)
    txbox(slide, ptitle, px+Inches(0.18), py+Inches(0.14), pw-Inches(0.36), Inches(0.38),
          font_name=SERIF, font_size=Pt(14), bold=True, color=TEXT_MAIN)
    txbox(slide, psub, px+Inches(0.18), py+Inches(0.5), pw-Inches(0.36), Inches(0.28),
          font_name=SANS, font_size=Pt(9.5), color=TEXT_MUTED)
    # Panel-specific content
    if pi == 0:  # Alerts
        for ai2, (sev, stxt, sdesc) in enumerate([
            ("CRITICAL", "Seat Conflict — Middle R1 S3", "Two dignitaries assigned to same seat."),
            ("WATCH",    "H.E. Rev. Egyir — Pending",    "Service started 20 mins ago. No arrival."),
        ]):
            acol = RED_ACC if sev == "CRITICAL" else AMBER_ACC
            abg = RGBColor(0x18,0x06,0x06) if sev == "CRITICAL" else RGBColor(0x18,0x14,0x06)
            ay3 = py + Inches(0.9) + ai2 * Inches(0.82)
            rounded_rect(slide, px+Inches(0.18), ay3, pw-Inches(0.36), Inches(0.7),
                         fill_color=abg, line_color=acol, line_width=Pt(0.5))
            txbox(slide, sev, px+Inches(0.28), ay3+Inches(0.06), pw-Inches(0.56), Inches(0.24),
                  font_name=MONO, font_size=Pt(7), bold=True, color=acol)
            txbox(slide, stxt, px+Inches(0.28), ay3+Inches(0.26), pw-Inches(0.56), Inches(0.24),
                  font_name=SANS, font_size=Pt(10), bold=True, color=TEXT_MAIN)
            txbox(slide, sdesc, px+Inches(0.28), ay3+Inches(0.46), pw-Inches(0.56), Inches(0.22),
                  font_name=SANS, font_size=Pt(9), color=TEXT_MUTED)
    elif pi == 1:  # Feed
        feed_items = [
            ("✅","Archbishop Mensah — Seated","09:14 · Middle Section R1/S1"),
            ("🚶","Pastor Oti-Boateng — Arrived","09:08 · Checked in"),
            ("🪑","Bishop Kwarteng — Seated","09:03 · Left Section R1/S2"),
        ]
        for fi, (ficon, flabel, ftime) in enumerate(feed_items):
            fy2 = py + Inches(0.9) + fi * Inches(0.56)
            txbox(slide, ficon, px+Inches(0.18), fy2, Inches(0.36), Inches(0.4), font_size=Pt(14))
            txbox(slide, flabel, px+Inches(0.6), fy2, pw-Inches(0.78), Inches(0.28),
                  font_name=SANS, font_size=Pt(10), bold=True, color=TEXT_MAIN)
            txbox(slide, ftime, px+Inches(0.6), fy2+Inches(0.26), pw-Inches(0.78), Inches(0.24),
                  font_name=SANS, font_size=Pt(9), color=TEXT_FAINT)
    elif pi == 2:  # Seating occupancy
        sections_occ = [
            ("Left",  RED_SEC,    90),
            ("Middle",BLUE_SEC,   71),
            ("Right", YELLOW_SEC, 50),
        ]
        for si3, (slbl, scol, pct) in enumerate(sections_occ):
            sy3 = py + Inches(0.9) + si3 * Inches(0.56)
            dot3 = slide.shapes.add_shape(9, px+Inches(0.2), sy3+Inches(0.08),
                                          Inches(0.13), Inches(0.13))
            dot3.fill.solid(); dot3.fill.fore_color.rgb = scol
            dot3.line.fill.background(); dot3.shadow.inherit = False
            txbox(slide, slbl, px+Inches(0.4), sy3, Inches(0.85), Inches(0.3),
                  font_name=SANS, font_size=Pt(9), color=TEXT_MUTED)
            progress_bar(slide, px+Inches(1.3), sy3+Inches(0.08),
                         pw-Inches(1.65), pct, scol)
            txbox(slide, f"{pct}%", px+pw-Inches(0.54), sy3, Inches(0.46), Inches(0.3),
                  font_name=MONO, font_size=Pt(8), color=TEXT_MUTED, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — Role-Based Access
# ══════════════════════════════════════════════════════════════════════════
slide = add_slide()
fill_slide_bg(slide)
add_logo(slide)

eyebrow(slide, "Access Control", Inches(0.55), Inches(0.82))
heading(slide, "The Right Tools for Every Role", Inches(0.55), Inches(1.08), size=Pt(38))

body(slide,
     "Everyone on the protocol team has access to exactly what they need — nothing more, nothing less.",
     Inches(0.55), Inches(2.0), Inches(9), Inches(0.5), size=Pt(12))

roles2 = [
    ("👑", "Administrator", "Full Control",
     ["Create & delete conferences","Manage all sessions",
      "Edit any seating plan","Manage user accounts",
      "Access dignitary directory","View control centre"]),
    ("✏️", "Editor", "Content Manager",
     ["Create new conferences","Add & edit sessions",
      "Assign & move seats","Add dignitaries to sessions",
      "Update attendance status","View control centre"]),
    ("🎫", "Protocol Member", "Field Observer",
     ["View seating maps","See dignitary profiles",
      "Update their own assignments","View attendance status"]),
]
rw = Inches(4.0); rh = Inches(4.5); rgap = Inches(0.24)
for ri3, (icon5, role_name, role_sub, perms) in enumerate(roles2):
    rx3 = Inches(0.55) + ri3 * (rw + rgap)
    ry3 = Inches(2.68)
    bc = RGBColor(0x18,0x2E,0x20) if ri3 == 0 else BG_PANEL
    lc = GOLD_DIM if ri3 == 0 else LINE_SOFT
    rounded_rect(slide, rx3, ry3, rw, rh, fill_color=bc, line_color=lc)
    txbox(slide, icon5, rx3+Inches(0.18), ry3+Inches(0.18), Inches(0.6), Inches(0.55), font_size=Pt(28))
    txbox(slide, role_name.upper(), rx3+Inches(0.18), ry3+Inches(0.82), rw-Inches(0.36), Inches(0.26),
          font_name=MONO, font_size=Pt(8), color=GOLD_DIM, bold=True)
    txbox(slide, role_sub, rx3+Inches(0.18), ry3+Inches(1.08), rw-Inches(0.36), Inches(0.45),
          font_name=SERIF, font_size=Pt(20), bold=True, color=TEXT_MAIN)
    for pi2, perm in enumerate(perms):
        py3 = ry3 + Inches(1.65) + pi2 * Inches(0.44)
        txbox(slide, "→", rx3+Inches(0.18), py3, Inches(0.28), Inches(0.34),
              font_name=SANS, font_size=Pt(11), color=GOLD_DIM)
        txbox(slide, perm, rx3+Inches(0.5), py3, rw-Inches(0.68), Inches(0.34),
              font_name=SANS, font_size=Pt(11), color=TEXT_SOFT)

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — Getting Started
# ══════════════════════════════════════════════════════════════════════════
slide = add_slide()
fill_slide_bg(slide)
add_logo(slide)

eyebrow(slide, "Quick Start", Inches(0.55), Inches(0.82))
heading(slide, "Up & Running in Minutes", Inches(0.55), Inches(1.08), size=Pt(42))

body(slide,
     "Getting started takes four simple steps. Once your first conference is set up, future events take even less time.",
     Inches(0.55), Inches(2.0), Inches(9), Inches(0.6), size=Pt(12))

qs_steps = [
    ("1", "🔑", "Sign In",            "Open the app in your browser and sign in. New users register in seconds."),
    ("2", "🏟️", "Create a Conference","Name the event, set the date, pick the auditorium, and add sessions."),
    ("3", "👤", "Add Dignitaries",    "Build or reuse the dignitary directory, then assign to each session."),
    ("4", "🗺️", "Seat & Track",       "Use the floor plan to assign seats, then go live on event day."),
]
sw3 = Inches(2.9); sh3 = Inches(3.5); sgap3 = Inches(0.24)
for si4, (num2, icon6, title3, desc3) in enumerate(qs_steps):
    sx4 = Inches(0.55) + si4 * (sw3 + sgap3)
    sy4 = Inches(2.82)
    rounded_rect(slide, sx4, sy4, sw3, sh3, fill_color=BG_PANEL, line_color=LINE_SOFT)
    txbox(slide, num2, sx4+Inches(0.18), sy4+Inches(0.16), Inches(0.6), Inches(0.65),
          font_name=SERIF, font_size=Pt(36), bold=True, color=GOLD_DIM)
    txbox(slide, icon6, sx4+sw3-Inches(0.65), sy4+Inches(0.18), Inches(0.55), Inches(0.55),
          font_size=Pt(24))
    txbox(slide, title3, sx4+Inches(0.18), sy4+Inches(0.95), sw3-Inches(0.36), Inches(0.48),
          font_name=SERIF, font_size=Pt(17), bold=True, color=TEXT_MAIN)
    txbox(slide, desc3, sx4+Inches(0.18), sy4+Inches(1.52), sw3-Inches(0.36), Inches(1.7),
          font_name=SANS, font_size=Pt(11), color=TEXT_MUTED, wrap=True)
    # Arrow
    if si4 < 3:
        ax4 = sx4 + sw3 + Inches(0.02)
        ay4 = sy4 + sh3/2 - Inches(0.2)
        txbox(slide, "›", ax4, ay4, sgap3+Inches(0.04), Inches(0.4),
              font_size=Pt(22), color=TEXT_FAINT, align=PP_ALIGN.CENTER)

# Pro tip card
rounded_rect(slide, Inches(0.55), Inches(6.46), Inches(12.25), Inches(0.76),
             fill_color=BG_PANEL, line_color=LINE_SOFT)
txbox(slide, "💡", Inches(0.75), Inches(6.55), Inches(0.45), Inches(0.5), font_size=Pt(18))
txbox(slide,
      "Pro tip: Use Import Arrangement to copy an existing session's seating plan into a new session — saving hours of re-entry.",
      Inches(1.28), Inches(6.57), Inches(11.2), Inches(0.46),
      font_name=SANS, font_size=Pt(11), color=TEXT_SOFT, wrap=True)

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — Closing
# ══════════════════════════════════════════════════════════════════════════
slide = add_slide()
fill_slide_bg(slide)

# Central glow
for gsize, gcol in [(Inches(6), RGBColor(0x12,0x28,0x1A)), (Inches(4), RGBColor(0x0D,0x20,0x14))]:
    gx = (W - gsize) / 2; gy = (H - gsize) / 2
    gsh = slide.shapes.add_shape(9, gx, gy, gsize, gsize)
    gsh.fill.solid(); gsh.fill.fore_color.rgb = gcol
    gsh.line.fill.background(); gsh.shadow.inherit = False

# Logo centred
if os.path.exists(LOGO_PATH):
    lw2, lh2 = Inches(2.8), Inches(0.9)
    lx2 = (W - lw2) / 2; ly2 = Inches(1.0)
    slide.shapes.add_picture(LOGO_PATH, lx2, ly2, lw2, lh2)

# Gold rule
horizontal_rule(slide, Inches(3.8), Inches(2.1), Inches(5.75), color=GOLD_DIM)

# Main close heading
txbox(slide, "Dignitary Protocol,",
      Inches(1.0), Inches(2.3), Inches(11.33), Inches(1.05),
      font_name=SERIF, font_size=Pt(56), bold=True,
      color=GOLD_BRIGHT, align=PP_ALIGN.CENTER)
txbox(slide, "Perfected.",
      Inches(1.0), Inches(3.2), Inches(11.33), Inches(1.05),
      font_name=SERIF, font_size=Pt(56), bold=True,
      color=GOLD_MAIN, align=PP_ALIGN.CENTER)

txbox(slide,
      "Pastors' Protocol Central brings peace of mind to every conference\n"
      "— so your team can focus on honouring your guests, not managing spreadsheets.",
      Inches(2.0), Inches(4.4), Inches(9.33), Inches(1.0),
      font_name=SANS, font_size=Pt(14), color=TEXT_MUTED,
      align=PP_ALIGN.CENTER, wrap=True)

horizontal_rule(slide, Inches(3.8), Inches(5.55), Inches(5.75), color=GOLD_DIM)

# Pills
badge_pill(slide, "✝ God's Love Tabernacle",
           Inches(4.0), Inches(5.78),
           BG_PANEL, GOLD_DIM, GOLD_MAIN, font_size=Pt(10))
badge_pill(slide, "Pastors' Protocol Central v1",
           Inches(4.0)+Inches(1.6), Inches(5.78),
           BG_PANEL, LINE_SOFT, TEXT_SOFT, font_size=Pt(10))

# ══════════════════════════════════════════════════════════════════════════
#  Save
# ══════════════════════════════════════════════════════════════════════════
out = os.path.join(os.path.dirname(__file__), "Pastors_Protocol_How_It_Works.pptx")
prs.save(out)
print(f"✓ Saved: {out}")
